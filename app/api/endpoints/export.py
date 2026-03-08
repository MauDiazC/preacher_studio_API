from fastapi import APIRouter, Depends
from fastapi.responses import StreamingResponse
from app.core.security import get_current_user
from app.repository.sermon_repository import sermon_repo
from app.core.exceptions import EntityNotFoundException
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import io

router = APIRouter(prefix="/export", tags=["Exportación"])


@router.get("/{sermon_id}/pdf", summary="Exportar sermón a PDF")
async def export_to_pdf(sermon_id: str, user_id: str = Depends(get_current_user)):
    res = sermon_repo.get_by_id(sermon_id, user_id)
    if not res.data:
        raise EntityNotFoundException(message="Sermón no encontrado para exportar.")

    sermon = res.data
    buffer = io.BytesIO()
    p = canvas.Canvas(buffer, pagesize=letter)

    # Generación simple de PDF
    p.setFont("Helvetica-Bold", 18)
    p.drawString(50, 750, sermon['title'])

    p.setFont("Helvetica-Oblique", 12)
    p.drawString(50, 730, f"Pasaje: {sermon.get('main_passage') or 'N/A'}")

    p.setFont("Helvetica", 11)
    text_object = p.beginText(50, 700)
    content = sermon.get("content", "")
    for line in content.split("\n"):
        text_object.textLine(line)
    p.drawText(text_object)

    p.showPage()
    p.save()

    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename=sermon_{sermon_id}.pdf"},
    )


@router.get("/{sermon_id}/pptx", summary="Exportar sermón a PowerPoint (Compatible con Keynote)")
async def export_to_pptx(sermon_id: str, user_id: str = Depends(get_current_user)):
    res = sermon_repo.get_by_id(sermon_id, user_id)
    if not res.data:
        raise EntityNotFoundException(message="Sermón no encontrado para exportar.")

    sermon = res.data
    prs = Presentation()
    
    # Diapositiva de Título
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = sermon["title"]
    subtitle.text = f"Pasaje: {sermon.get('main_passage') or 'N/A'}"

    # Diapositivas de Contenido (dividir por párrafos o puntos)
    content = sermon.get("content", "")
    paragraphs = [p for p in content.split("\n") if p.strip()]
    
    bullet_slide_layout = prs.slide_layouts[1]
    for p_text in paragraphs:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = sermon["title"]
        tf = body_shape.text_frame
        tf.text = p_text

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename=sermon_{sermon_id}.pptx"
        },
    )
